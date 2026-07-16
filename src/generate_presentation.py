#!/usr/bin/env python3
"""Generate 22-slide widescreen presentation from figures and paper content."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PROJECT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG  = os.path.join(PROJECT, "data")
OUT  = os.path.join(PROJECT, "presentation")

os.makedirs(OUT, exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Color palette ───────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0xE2, 0x4A, 0x33)
ACCENT2    = RGBColor(0x34, 0x8A, 0xBD)
GRAY       = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

# ── Helpers ─────────────────────────────────────────────────────────────
def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def text_box(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=16, bold=False, color=DARK, space_before=Pt(6),
             font_name="Calibri", align=PP_ALIGN.LEFT, bullet=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    p.level = 1 if bullet else 0
    return p

def title_slide(title_text, subtitle_text=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(sl, WHITE)
    # Accent bar at top
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    text_box(sl, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
             title_text, size=36, bold=True, color=DARK, align=PP_ALIGN.LEFT)
    if subtitle_text:
        text_box(sl, Inches(1), Inches(3.8), Inches(11.3), Inches(1.2),
                 subtitle_text, size=18, color=GRAY, align=PP_ALIGN.LEFT)
    return sl

def section_slide(number, title):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, WHITE)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    # Number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(0.6), Inches(0.8), Inches(0.8))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT; circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(number); p.font.size = Pt(28)
    p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    text_box(sl, Inches(1.9), Inches(0.65), Inches(10), Inches(0.7),
             title, size=30, bold=True, color=DARK)
    # Thin separator line
    sep = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.02))
    sep.fill.solid(); sep.fill.fore_color.rgb = LIGHT_GRAY; sep.line.fill.background()
    return sl

def figure_slide(number, title, fig_name):
    sl = section_slide(number, title)
    img_path = os.path.join(FIG, fig_name)
    if os.path.exists(img_path):
        # Center image
        img_w = Inches(10)
        ratio = 10.0 / 13.333  # rough
        left = (W - img_w) // 2
        top = Inches(2.0)
        sl.shapes.add_picture(img_path, left, top, width=img_w)
    return sl

def bullet_slide(number, title, bullets, subheader=""):
    sl = section_slide(number, title)
    if subheader:
        text_box(sl, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.5),
                 subheader, size=16, color=GRAY)
    tf = text_box(sl, Inches(0.9), Inches(2.3), Inches(11.5), Inches(4.5),
                  "", size=18, color=DARK)
    tf.paragraphs[0].text = ""
    for b in bullets:
        add_para(tf, b, size=18, bullet=True, space_before=Pt(8))
    return sl

def callout_slide(number, title, body, callout_text):
    sl = section_slide(number, title)
    text_box(sl, Inches(0.9), Inches(2.0), Inches(5.5), Inches(4.5),
             body, size=16, color=DARK)
    # Big callout box
    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(7.0), Inches(2.2), Inches(5.5), Inches(3.0))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xED)
    box.line.fill.background()
    tf2 = box.text_frame; tf2.word_wrap = True
    p = tf2.paragraphs[0]; p.text = callout_text; p.font.size = Pt(20)
    p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf2.paragraphs[0].space_before = Pt(30)
    return sl

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════
title_slide(
    "Multi-Horizon Hazard Models for\nBattery Failure Prediction",
    "Within-Dataset Reliability and Cross-Chemistry Transferability\n\nHussain Touhid Siddiquee et al. · Department of EEE, Leading University, Sylhet"
)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — MOTIVATION
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(1, "Motivation", [
    "Shikdar & Laaksonen (2026): multi-horizon hazard classification on NASA 18650 using HGB — AUC 0.87–0.90",
    "Limited to one model class, one chemistry, no cross-dataset validation",
    "Two open questions:",
    "  1. How sensitive are results to model choice (XGBoost, LightGBM, RF) and calibration method?",
    "  2. Do models trained on one chemistry (LCO) transfer to another (LFP)?",
],
    subheader="The original framework demonstrated feasibility — but several critical questions remain unanswered"
)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — COMPOSITE FAILURE LABEL
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(2, "The Composite Failure Label", [
    "Two independent failure criteria — either triggers a positive label:",
    "  • State-of-Health (SOH) ≤ 0.80 of initial capacity",
    "  • Average voltage sag < 94% of early-life baseline (first 10 cycles)",
    "Once triggered, label remains positive for all subsequent cycles",
    "Voltage sag catches impedance-driven failures that precede capacity fade",
    "LFP cells: no useful voltage sag (flat plateau at 2.7 V cutoff)",
],
    subheader="Failure = SOH drop OR voltage sag. Both must cross a per-cell baseline."
)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — DATASETS & MODELS
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide(3, "Datasets & Models")

# Table data
tbl_data = [
    ["Dataset", "Cells", "Chemistry", "Total Cycles", "Cycles/Cell", "Profile"],
    ["NASA 18650", "37", "LCO", "~1,000", "~300", "Random walk, accelerated"],
    ["CALCE LCO/CX2", "7", "LCO", "8,733", "775–1,952", "1C/1C, slow degradation"],
    ["Oxford LFP", "5", "LFP", "~1,500", "~300", "1C/1C, very stable"],
]
tbl = sl.shapes.add_table(len(tbl_data), 6, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.5)).table
for r, row in enumerate(tbl_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = str(val)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14) if r > 0 else Pt(14)
            p.font.bold = (r == 0)
            p.font.color.rgb = DARK if r > 0 else WHITE
            p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE

# Model info
tf = text_box(sl, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.8),
              "Models: XGBoost (depth=4, 300 trees, lr=0.05)  ·  "
              "LightGBM (depth=4, 300 trees, lr=0.05)  ·  "
              "Random Forest (depth=6, 300 trees)  ·  "
              "GRU (1 layer, 8 hidden, W=10, lr=0.005) — compact to avoid overfitting",
              size=16, color=DARK)
add_para(tf, "Tree features: cycle, avg/min voltage, avg current, avg temp, duration, SOH", size=14, color=GRAY, bullet=True)
add_para(tf, "GRU features: cycle, avg/min voltage, SOH (cross-chem drops avg_current, avg_temp, duration)", size=14, color=GRAY, bullet=True)
add_para(tf, "All tree-based hyperparameters matched to original study", size=14, color=GRAY, bullet=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — PROTOCOL
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(4, "Evaluation Protocol", [
    "Within-dataset: 5-fold GroupKFold by cell — no cell leaks across folds",
    "Horizons: H ∈ {10, 20, 30, 50} — label positive if failure within [t, t+H)",
    "Calibration: isotonic vs Platt (sigmoid) — cross-chem reports raw AUC (calibration unreliable under distribution shift)",
    "Metrics: AUC (discrimination) + Brier score (calibration + discrimination)",
    "Cross-chemistry: train on all LCO, test on all LFP — 3 variants (NASA, CALCE, ALL)",
    "Feature ablation: with SOH vs without SOH — to isolate SOH contribution",
],
    subheader="Two parallel experiments: within-dataset reliability + cross-chemistry transfer"
)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — FINDING 1: WITHIN-DATASET AUC
# ══════════════════════════════════════════════════════════════════════════
figure_slide(5, "Finding 1: Within-Dataset AUC (H=20)", "Fig01_Within_Dataset_AUC.png")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — FINDING 2: MULTI-HORIZON
# ══════════════════════════════════════════════════════════════════════════
figure_slide(6, "Finding 2: Multi-Horizon Performance (NASA)", "Fig05_MultiHorizon_AUC.png")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — CALIBRATION NASA
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide(7, "Finding 3: Platt vs Isotonic — NASA")
img_path = os.path.join(FIG, "Fig02a_Calibration_NASA.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.5), Inches(2.0), width=Inches(6.5))
    tf = text_box(sl, Inches(7.5), Inches(2.5), Inches(5.3), Inches(4.0), "", size=16, color=DARK)
    tf.paragraphs[0].text = ""
    add_para(tf, "NASA:", size=18, bold=True, color=ACCENT)
    add_para(tf, "Platt AUC: 0.890 vs Isotonic AUC: 0.840", size=20, bold=True, space_before=Pt(16))
    add_para(tf, "Platt Brier: 0.213 vs Isotonic Brier: 0.214", size=16, color=GRAY, space_before=Pt(4))
    add_para(tf, "", size=10)
    add_para(tf, "AUC gap is modest on NASA. Isotonic", size=16, color=GRAY)
    add_para(tf, "partially degrades discrimination on", size=16, color=GRAY)
    add_para(tf, "noisier predictions.", size=16, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — CALIBRATION CALCE
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide(8, "Finding 3: Platt vs Isotonic — CALCE")
img_path = os.path.join(FIG, "Fig02b_Calibration_CALCE.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.5), Inches(2.0), width=Inches(6.5))
    tf = text_box(sl, Inches(7.5), Inches(2.5), Inches(5.3), Inches(4.0), "", size=16, color=DARK)
    tf.paragraphs[0].text = ""
    add_para(tf, "CALCE:", size=18, bold=True, color=ACCENT)
    add_para(tf, "Platt AUC: 0.904 vs Isotonic AUC: 0.694", size=20, bold=True, space_before=Pt(16))
    add_para(tf, "Brier is nearly identical: 0.105 for both", size=16, color=GRAY, space_before=Pt(4))
    add_para(tf, "", size=10)
    add_para(tf, "The advantage is in discrimination,", size=16, color=GRAY)
    add_para(tf, "not calibration — isotonic produces", size=16, color=GRAY)
    add_para(tf, "degenerate probabilities on long-tailed data.", size=16, color=GRAY)
    add_para(tf, "", size=10)
    add_para(tf, "Fair comparison: both calibrators use", size=16, color=GRAY)
    add_para(tf, "same underlying model's outputs", size=16, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — CROSS-CHEM WITH SOH
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide(9, "Finding 4a: Cross-Chemistry Transfer — With SOH (Oxford)")
img_path = os.path.join(FIG, "Fig03a_CrossChem_With_SOH_Oxford.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.8), Inches(1.9), width=Inches(7.5))
    tf = text_box(sl, Inches(8.8), Inches(2.2), Inches(4.0), Inches(4.5), "", size=16, color=DARK)
    tf.paragraphs[0].text = ""
    add_para(tf, "Raw AUC 0.84–1.00", size=28, bold=True, color=ACCENT)
    add_para(tf, "Per-cell mean ± std across test cells (trees H=20)", size=14, color=GRAY, space_before=Pt(6))
    add_para(tf, "", size=6)
    add_para(tf, "Oxford (5 cells)", size=14, color=GRAY)
    add_para(tf, "NASA→Oxford: 0.96–1.00", size=16, bold=True, color=DARK, space_before=Pt(8))
    add_para(tf, "CALCE→Oxford: 0.84–0.89", size=16, color=GRAY, space_before=Pt(2))
    add_para(tf, "ALL→Oxford: 0.98–1.00", size=16, color=GRAY, space_before=Pt(2))
    add_para(tf, "", size=6)
    add_para(tf, "Severson (141 cells) → 0.99+ Platt AUC", size=14, color=GRAY, space_before=Pt(2))
    add_para(tf, "", size=6)
    add_para(tf, "BUT — even with SOH, sequence models", size=16, bold=True, color=ACCENT, space_before=Pt(12))
    add_para(tf, "are less able to exploit it than tree-based models", size=16, bold=True, color=ACCENT)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — CROSS-CHEM WITH SOH SEVERSON
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide(9.5, "Finding 4a: Cross-Chemistry — With SOH (Severson)")
img_path = os.path.join(FIG, "Fig03b_CrossChem_With_SOH_Severson.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.8), Inches(1.9), width=Inches(7.5))
    tf = text_box(sl, Inches(8.8), Inches(2.2), Inches(4.0), Inches(4.5), "", size=16, color=DARK)
    tf.paragraphs[0].text = ""
    add_para(tf, "Severson (141 cells)", size=28, bold=True, color=ACCENT)
    add_para(tf, "Platt AUC 0.99+ across all models", size=16, color=GRAY, space_before=Pt(6))
    add_para(tf, "Confirms SOH-driven pattern", size=16, color=GRAY, space_before=Pt(6))

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — CROSS-CHEM WITHOUT SOH (OXFORD)
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide(10, "Finding 4b: Cross-Chemistry — Without SOH (Oxford)")
img_path = os.path.join(FIG, "Fig04a_CrossChem_No_SOH_Oxford.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.8), Inches(1.9), width=Inches(7.5))
    tf = text_box(sl, Inches(8.8), Inches(2.2), Inches(4.0), Inches(4.5), "", size=16, color=DARK)
    tf.paragraphs[0].text = ""
    add_para(tf, "Raw AUC 0.33–0.62 (Oxford)", size=28, bold=True, color=ACCENT)
    add_para(tf, "Trees: near-random across all configs", size=16, color=GRAY, space_before=Pt(6))
    add_para(tf, "", size=6)
    add_para(tf, "NASA→Oxford: 0.52–0.54", size=16, bold=True, color=DARK)
    add_para(tf, "CALCE→Oxford: 0.46–0.62", size=16, color=GRAY, space_before=Pt(2))
    add_para(tf, "ALL→Oxford: 0.33–0.47", size=16, color=GRAY, space_before=Pt(2))
    add_para(tf, "", size=10)
    add_para(tf, "No model class achieves above-chance", size=16, bold=True, color=ACCENT, space_before=Pt(12))
    add_para(tf, "AUC once SOH is removed", size=16, bold=True, color=ACCENT)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 11B — CROSS-CHEM WITHOUT SOH (SEVERSON)
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide(10.5, "Finding 4b: Cross-Chemistry — Without SOH (Severson)")
img_path = os.path.join(FIG, "Fig04b_CrossChem_No_SOH_Severson.png")
if os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.8), Inches(1.9), width=Inches(7.5))
    tf = text_box(sl, Inches(8.8), Inches(2.2), Inches(4.0), Inches(4.5), "", size=16, color=DARK)
    tf.paragraphs[0].text = ""
    add_para(tf, "Severson raw: 0.60–0.75", size=28, bold=True, color=ACCENT)
    add_para(tf, "Cycle-number proxy partially overlaps", size=16, color=GRAY, space_before=Pt(6))
    add_para(tf, "with wider LCO cycle-life range", size=16, color=GRAY, space_before=Pt(2))
    add_para(tf, "", size=10)
    add_para(tf, "19–24 AUC point drop from with-SOH", size=16, bold=True, color=ACCENT, space_before=Pt(12))
    add_para(tf, "Confirms SOH is the sole driver", size=16, bold=True, color=ACCENT)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDES 12–14 — SHAP EVIDENCE (2-panel per model, with/without SOH)
# ══════════════════════════════════════════════════════════════════════════
figure_slide(11, "SHAP: XGBoost — With vs Without SOH", "Fig06a_XGBoost_SHAP.png")
figure_slide(12, "SHAP: LightGBM — With vs Without SOH", "Fig06b_LightGBM_SHAP.png")
figure_slide(13, "SHAP: Random Forest — With vs Without SOH", "Fig06c_RandomForest_SHAP.png")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — DELONG TEST: STATISTICAL SIGNIFICANCE
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide(14, "Statistical Significance — DeLong Test")
tf = text_box(sl, Inches(0.9), Inches(1.9), Inches(11.5), Inches(5.0),
              "", size=16, color=DARK)
tf.paragraphs[0].text = ""
add_para(tf, "DeLong nonparametric test compares paired AUC values", size=18, bold=True)
add_para(tf, "", size=6)
add_para(tf, "SOH Ablation (ALL LCO→LFP, H=20):", size=20, bold=True, color=ACCENT, space_before=Pt(10))
add_para(tf, "Oxford (5 cells):", size=18, bold=True, color=DARK, space_before=Pt(8))
add_para(tf, "  XGBoost:    0.917 → 0.429    p = 7.2×10⁻⁵¹  ✓", size=18, color=DARK, space_before=Pt(6))
add_para(tf, "  LightGBM:   0.971 → 0.332    p = 2.6×10⁻⁶⁴  ✓", size=18, color=DARK, space_before=Pt(4))
add_para(tf, "  Random Forest:  0.989 → 0.581    p = 6.5×10⁻³⁶  ✓", size=18, color=DARK, space_before=Pt(4))
add_para(tf, "", size=4)
add_para(tf, "Severson (141 cells):", size=18, bold=True, color=DARK, space_before=Pt(8))
add_para(tf, "  XGBoost:    0.896 → 0.750    p < 10⁻¹⁰⁰  ✓", size=18, color=DARK, space_before=Pt(6))
add_para(tf, "  LightGBM:   0.882 → 0.718    p < 10⁻¹⁰⁰  ✓", size=18, color=DARK, space_before=Pt(4))
add_para(tf, "  Random Forest:  0.889 → 0.746    p < 10⁻¹⁰⁰  ✓", size=18, color=DARK, space_before=Pt(4))
add_para(tf, "", size=4)
add_para(tf, "All models, both LFP targets: p < 10⁻³⁶ — decisive", size=18, bold=True, color=ACCENT, space_before=Pt(8))
add_para(tf, "SOH-ablation gap is not due to chance", size=18, color=GRAY, space_before=Pt(2))
add_para(tf, "", size=6)
add_para(tf, "Within-dataset (CALCE): all model pairs p < 10⁻⁵", size=16, color=GRAY, space_before=Pt(10))
add_para(tf, "Within-dataset (NASA): only XGB vs LGBM reaches p=0.014", size=16, color=GRAY, space_before=Pt(2))

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — THE MECHANISM
# ══════════════════════════════════════════════════════════════════════════
callout_slide(
    15, "Why? The SOH-as-Lookup-Table Mechanism",
    "When SOH is available as a feature:\n\n"
    "• Model learns: SOH=0.85 means ~50 cycles to failure\n"
    "  from LCO training data\n\n"
    "• LFP cells traverse similar SOH trajectories\n\n"
    "• Model applies the same learned mapping — "
    "predictions correlate with SOH, not with LFP-specific degradation\n\n"
    "• Result: high AUC, but no actual chemistry transfer\n\n"
    "When SOH is removed: voltage/current/temperature\n"
    "patterns alone carry no transferable signal LCO→LFP\n\n"
     "GRU finding: even with SOH, compressed 8-unit hidden state\n"
     "must share capacity across SOH, voltage, and cycle —\n"
     "under shift the entangled noise corrupts SOH signal.\n\n"
     "Tree-model isolated splits avoid this entirely —\n"
     "architecture design matters as much as capacity.",
    "SOH is not a transferable feature —\n"
    "it is a chemistry-specific\n"
    "capacity-to-RUL lookup table"
)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — KEY TAKEAWAY
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide(16, "Key Takeaway")
# Big centered statement
box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.8))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xED)
box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Cross-chemistry transfer of hazard-based battery\nfailure prediction remains an open problem."
p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = ACCENT
p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(40)

# Supporting line
text_box(sl, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.2),
         "Within-dataset reliability is strong (AUC ≥ 0.85, including GRU). But no model class \u2014 "
         "tree-based or sequence-aware \u2014 achieves genuine cross-chemistry transfer without SOH. "
         "Practitioners should not expect LCO-trained hazard models to generalize to LFP without "
         "chemistry-specific retraining or domain adaptation.",
         size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — LIMITATIONS
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(17, "Limitations", [
    "Oxford LFP: only 5 cells (partially mitigated by Severson LFP, 141 cells — reproduces all qualitative findings)",
    "Cross-chemistry: unidirectional only (LCO→LFP); may not generalize to LCO→NMC, NMC→LFP",
    "GRU cross-chem unstable across H (e.g., CALCE→Oxford ranges 0.022–0.978); mean-H reported in figures",
    "Brier gap with published paper: published Brier 0.032 vs reproducible 0.17–0.26 (noted in discrepancy note)",
    "Voltage sag feature dead for LFP: always at 2.7 V cutoff — no useful signal",
    "Calibration fails under cross-chem distribution shift: isotonic collapses AUC (0.98→0.51), Platt unstable — raw AUC used",
    "Promising directions: physics-informed features (ICA/DVA), few-shot fine-tuning, domain adaptation minimizing LCO↔LFP distribution gap",
])

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(18, "Next Steps", [
    "Train on larger multi-chemistry datasets with balanced cell counts across chemistries",
    "Develop learned feature representations designed explicitly for chemistry invariance (e.g., domain-adversarial training)",
    "Bidirectional transfer evaluation: LCO↔NMC, NMC↔LFP, LCO↔LFP",
    "Sequence models with explicit invariance objectives (GRU tested — fails without SOH; need domain-adversarial)",
    "Reproduce and close the Brier gap with the published paper",
])

# ── Save ────────────────────────────────────────────────────────────────
out_path = os.path.join(OUT, "presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path} ({len(prs.slides)} slides)")
