#!/usr/bin/env python3
"""Generate simplified 21-slide presentation with updated figures and GRU data."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PROJECT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG  = os.path.join(PROJECT, "data")
OUT  = os.path.join(PROJECT, "presentation")

os.makedirs(OUT, exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Colors ──
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0xE2, 0x4A, 0x33)
GRAY       = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def text_box(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return tf

def add_para(tf, text, size=16, bold=False, color=DARK, space_before=Pt(6),
             align=PP_ALIGN.LEFT, bullet=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    p.space_before = space_before
    p.level = 1 if bullet else 0
    return p

def title_slide_big(title_text, subtitle_text=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, WHITE)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    text_box(sl, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
             title_text, size=36, bold=True, color=DARK)
    if subtitle_text:
        text_box(sl, Inches(1), Inches(3.8), Inches(11.3), Inches(1.2),
                 subtitle_text, size=18, color=GRAY)

def section_slide_simple(number, title):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl, WHITE)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    text_box(sl, Inches(0.7), Inches(0.3), Inches(11.9), Inches(0.7),
             title, size=28, bold=True, color=DARK)
    sep = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.1), Inches(11.9), Inches(0.01))
    sep.fill.solid(); sep.fill.fore_color.rgb = LIGHT_GRAY; sep.line.fill.background()
    return sl

def bullet_slide(number, title, bullets):
    sl = section_slide_simple(number, title)
    tf = text_box(sl, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5),
                  "", size=18, color=DARK)
    tf.paragraphs[0].text = ""
    for b in bullets:
        add_para(tf, b, size=18, bullet=True, space_before=Pt(8))
    return sl

def figure_slide(number, title, fig_name, caption=""):
    sl = section_slide_simple(number, title)
    img_path = os.path.join(FIG, fig_name)
    if os.path.exists(img_path):
        sl.shapes.add_picture(img_path, Inches(1.2), Inches(1.4), width=Inches(10.5))
    if caption:
        text_box(sl, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.8),
                 caption, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    return sl

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════
title_slide_big(
    "Multi-Horizon Hazard Models for\nBattery Failure Prediction",
    "Can we predict battery failure across different chemistries?\n\n"
    "Hussain Touhid Siddiquee et al. \u00b7 Department of EEE, Leading University, Sylhet"
)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(1, "The Problem", [
    "Batteries fail over time \u2014 we want to predict failure before it happens",
    "This helps EVs, phones, and grid storage operate safely",
    "Previous work tested only ONE battery type (LCO laptop batteries)",
    "Question 1: Do different AI models give different results?",
    "Question 2: Can we train on one chemistry and predict on another?",
])

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — WHAT DOES FAILURE MEAN
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(2, "What Does \u201cFailure\u201d Mean?", [
    "We check two vital signs, like a doctor checking health:",
    "1. State-of-Health (SOH) \u2014 battery capacity dropped below 80%",
    "2. Voltage Sag \u2014 power drops suddenly before capacity fades",
    "Either one triggered = \u201cfailure\u201d label",
    "Once triggered, stays triggered \u2014 batteries don\u2019t heal",
])

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — THE BATTERIES WE TESTED
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(3, "The Batteries We Tested", [
    "NASA 18650: 37 LCO cells (laptop chemistry), ~300 cycles each",
    "CALCE LCO/CX2: 7 LCO cells, 775\u20131952 cycles each",
    "Oxford LFP: 5 LFP cells (EV chemistry), ~300 cycles each",
    "LCO \u2192 train. LFP \u2192 test. Train on LCO, predict on LFP? That\u2019s the key question.",
])

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — HOW WE TESTED
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(4, "How We Tested", [
    "Four AI models: XGBoost, LightGBM, Random Forest, GRU (sequence model)",
    "GRU: 1 layer, 8 hidden units (compact to avoid overfitting), sliding window of 10 cycles per cell",
    "Four prediction windows: 10, 20, 30, 50 cycles ahead",
    "Two calibration methods: isotonic (freehand) vs Platt (smooth curve)",
    "\u26a0 Critical: test with SOH feature AND without it",
])

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — FINDING 1: WITHIN-DATASET
# ══════════════════════════════════════════════════════════════════════════
figure_slide(5, "Finding 1: Works on Known Batteries",
             "Fig01_Within_Dataset_AUC.png",
             "All models (including GRU) score AUC \u2265 0.85 on NASA and CALCE \u2014 the framework is reliable on familiar chemistries")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — FINDING 2: MULTI-HORIZON
# ══════════════════════════════════════════════════════════════════════════
figure_slide(6, "Finding 2: Longer Windows = Easier",
             "Fig05_MultiHorizon_AUC.png",
             "Predicting 50 cycles ahead is easier than 10 \u2014 more time gives more data to work with")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — FINDING 3: CALIBRATION
# ══════════════════════════════════════════════════════════════════════════
figure_slide(7, "Finding 3: Platt vs Isotonic — NASA",
             "Fig02a_Calibration_NASA.png",
             "NASA: Platt AUC=0.890 vs Isotonic AUC=0.840. Platt wins, but gap is modest.")
figure_slide(8, "Finding 3: Platt vs Isotonic — CALCE",
             "Fig02b_Calibration_CALCE.png",
             "CALCE: Platt AUC=0.904 vs Isotonic AUC=0.694. Dramatic gap on long-tailed data.")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — WHY PLATT IS BETTER
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(8, "Why Platt Calibration is Better", [
    "Think of isotonic as drawing a freehand squiggle to fit the data:",
    "\u2022 Follows every bump \u2014 works fine on small, clean datasets",
    "\u2022 But overfits on messy, long-tailed data (like CALCE\u2019s 8733 cycles)",
    "Platt uses a smooth S-curve instead:",
    "\u2022 Like taking a step back and seeing the big picture",
    "\u2022 More robust \u2014 it doesn\u2019t chase noise",
    "CALCE: Platt AUC=0.904 vs isotonic AUC=0.694 (discrimination gain); Brier nearly identical at 0.105",
])

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — FINDING 4a: WITH SOH
# ══════════════════════════════════════════════════════════════════════════
figure_slide(9, "Finding 4a: With SOH — Oxford",
             "Fig03a_CrossChem_With_SOH_Oxford.png",
             "Oxford (5 cells): tree models 0.84–1.00 with SOH. CALCE-inclusive GRU reversal (AUC≈0.03–0.12) signals class-imbalance transfer failure.")
figure_slide(10, "Finding 4a: With SOH — Severson",
             "Fig03b_CrossChem_With_SOH_Severson.png",
             "Severson (141 cells): confirms SOH-driven pattern. Platt AUC 0.99+ for ALL LCO→Severson.")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 11–12 — FINDING 4b: WITHOUT SOH
# ══════════════════════════════════════════════════════════════════════════
figure_slide(11, "Finding 4b: Without SOH — Oxford",
             "Fig04a_CrossChem_No_SOH_Oxford.png",
             "Oxford: trees near-random (0.33–0.62) — no genuine chemistry transfer.")
figure_slide(12, "Finding 4b: Without SOH — Severson",
             "Fig04b_CrossChem_No_SOH_Severson.png",
             "Severson: trees 0.60–0.75 (cycle-number proxy), 19–24 AUC drop from with-SOH condition.")

# ══════════════════════════════════════════════════════════════════════════
#  SLIDES 13–15 — SHAP EVIDENCE (2-panel per model: with/without SOH)
# ══════════════════════════════════════════════════════════════════════════
figure_slide(13, "SHAP: XGBoost — With vs Without SOH",
             "Fig06a_XGBoost_SHAP.png",
             "SHAP: top panel with SOH (SOH dominates), bottom panel without SOH (all features collapse). NASA→Oxford, H=20.")
figure_slide(14, "SHAP: LightGBM — With vs Without SOH",
             "Fig06b_LightGBM_SHAP.png",
             "SHAP: top panel with SOH (SOH dominates), bottom panel without SOH (all features collapse). NASA→Oxford, H=20.")
figure_slide(15, "SHAP: Random Forest — With vs Without SOH",
             "Fig06c_RandomForest_SHAP.png",
             "SHAP: top panel with SOH (SOH dominates), bottom panel without SOH (all features collapse). NASA→Oxford, H=20."))

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — THE SOH LOOKUP TABLE
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(16, "The \u201cSOH Lookup Table\u201d", [
    "When the model has SOH, it learns this rule:",
    "\u201cSOH = 85% \u2192 about 50 more cycles before failure\u201d",
    "It memorizes this from LCO training data",
    "LFP batteries happen to have similar SOH values\u2026",
    "\u2026so the model applies the same wrong rule to LFP",
    "Result: predictions look good (high AUC)",
    "But: the model learned nothing about LFP itself \u2014 just reused a lookup table",
    "GRU reveals architecture-specific fragility: 8-unit hidden state shares capacity across SOH + voltage + cycle → under shift, entangled noise corrupts SOH signal. AUC swing 0.011–0.986 within a single config — the instability itself is the finding.",
])

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — WHAT THIS MEANS
# ══════════════════════════════════════════════════════════════════════════
sl = section_slide_simple(17, "What This Means")
box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(1.5), Inches(2.2), Inches(10.3), Inches(2.5))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xED)
box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The hazard framework works reliably\nwithin the same battery chemistry\n(AUC \u2265 0.85, including GRU)."
p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = ACCENT
p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(20)

box2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.8))
box2.fill.solid(); box2.fill.fore_color.rgb = WHITE
box2.line.color.rgb = ACCENT; box2.line.width = Pt(2)
tf2 = box2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = ("But cross-chemistry transfer remains an open problem for all model classes.\n"
           "You cannot train on LCO and deploy on LFP without SOH \u2014 and SOH is a chemistry-specific lookup key.")
p2.font.size = Pt(16); p2.font.color.rgb = DARK
p2.font.name = "Calibri"; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(10)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — CAVEATS
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(18, "Caveats", [
    "Oxford LFP: only 5 cells \u2014 mitigated by Severson (141 cells) but both LFP only",
    "Severson LFP: fast-charging protocol (4C discharge), voltage sag uniformly uninformative across both LFP datasets",
    "Tested one direction only (LCO \u2192 LFP)",
    "GRU cross-chem instability (AUC range 0.011–0.986 within one config) from distributed hidden-state entanglement under covariate shift — an informative negative result showing GRUs poorly suited to cross-chem in single-seed setting",
    "Voltage sag useless for LFP (flat voltage plateau)",
    "Published Brier scores could not be reproduced (code mismatch)",
    "Calibration fails under cross-chem distribution shift: isotonic collapses AUC (0.98→0.51); raw AUC used instead",
    "SHAP confirms SOH is the sole driver of cross-chem AUC: each SHAP figure shows with/without SOH side by side",
    "Promising: physics-informed features (ICA/DVA), few-shot fine-tuning, domain adaptation minimizing LCO↔LFP gap",
])

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════
bullet_slide(19, "Next Steps", [
    "Test on larger datasets with more battery types",
    "Design features that are \u201cchemistry-agnostic\u201d",
    "GRU single-seed instability is itself the finding: distributed representations fail under chemistry shift. Multi-seed analysis + domain-adversarial training deferred",
    "Test both directions: LCO\u2192LFP AND LFP\u2192LCO",
])

# ── Save ──
out_path = os.path.join(OUT, "presentation_simple.pptx")
prs.save(out_path)
print(f"Saved: {out_path} ({len(prs.slides)} slides)")
